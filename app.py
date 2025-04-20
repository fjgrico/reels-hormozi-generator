import streamlit as st
import openai
import io
from docx import Document
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from prompts.generador_guiones import generar_guion

# Configuración de la página
st.set_page_config(page_title='Generador de Reels IA', page_icon='🎬', layout='centered')
st.title('🎬 Generador de Guiones estilo Hormozi con IA')

st.subheader('1. Genera tu guion brutal:')

with st.form('guion_form'):
    col1, col2 = st.columns(2)
    with col1:
        sector = st.text_input('¿A qué te dedicas?')
        producto = st.text_input('¿Qué vendes?')
        beneficios = st.text_area('Beneficios del producto/servicio')
    with col2:
        problemas = st.text_area('Problemas que soluciona')
        avatar = st.text_input('Describe tu cliente ideal')
        tematica = st.selectbox('Temática del Reel', ['Objeciones', 'Resultado', 'Historia', 'CTA'])

    tono = st.selectbox('Tono', ['Automático', 'Agresivo', 'Emocional', 'Inspirador', 'Polémico'])
    copy = st.selectbox('Fórmula', ['PAS', 'AIDA', 'FAB', 'Hormozi Style', 'Storytelling'])
    duracion = st.slider('Duración', 20, 60, 30)
    submit = st.form_submit_button('🚀 Generar Guion')

if submit:
    datos = {
        'a_que_se_dedica': sector,
        'producto_servicio': producto,
        'beneficios': beneficios,
        'problemas': problemas,
        'avatar': avatar,
        'tematica': tematica,
        'tono': tono,
        'copywriting': copy,
        'duracion': duracion
    }
    with st.spinner('Generando...'):
        resultado = generar_guion(datos)
        st.success('✅ Guion generado:')
        st.markdown(f"```{resultado}```")

        # DOCX
        def crear_docx(texto):
            doc = Document()
            doc.add_heading("Guion Estilo Hormozi", level=1)
            for linea in texto.split('\n'):
                doc.add_paragraph(linea)
            buffer = io.BytesIO()
            doc.save(buffer)
            return buffer.getvalue()

        # PDF (solución corregida)
        def crear_pdf(texto):
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            for linea in texto.split('\n'):
                pdf.multi_cell(0, 10, linea)
            buffer = io.BytesIO()
            pdf_bytes = pdf.output(dest='S').encode('latin1')
            buffer.write(pdf_bytes)
            buffer.seek(0)
            return buffer.read()

        # PPTX
        def crear_pptx(texto):
            prs = Presentation()
            layout = prs.slide_layouts[5]
            for bloque in texto.split('\n\n'):
                slide = prs.slides.add_slide(layout)
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(5))
                tf = txBox.text_frame
                tf.text = bloque
                for p in tf.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(24)
            buffer = io.BytesIO()
            prs.save(buffer)
            return buffer.getvalue()

        # BOTONES DE DESCARGA
        st.download_button("📄 Descargar como Word (.docx)", crear_docx(resultado), file_name="guion.docx")
        st.download_button("🧾 Descargar como PDF (.pdf)", crear_pdf(resultado), file_name="guion.pdf")
        st.download_button("📽️ Descargar como PowerPoint (.pptx)", crear_pptx(resultado), file_name="guion.pptx")

st.divider()
st.subheader('2. Próximamente: voz y vídeo IA (activado en tu versión pro)')
